import os
import tempfile
import subprocess
import shutil
import zipfile
import logging
from flask import Flask, request, send_file, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
import time

# Try to import pdf2image
try:
    import pdf2image
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logging.warning("pdf2image is not installed. PDF to JPG conversion will not be available.")

app = Flask(__name__, static_folder='.')
CORS(app)  # Enable CORS for all routes

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# LibreOffice Docker container details
LIBREOFFICE_CONTAINER = "libreoffice"

# Create uploads directory if it doesn't exist
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    try:
        os.makedirs(UPLOAD_FOLDER)
    except FileExistsError:
        pass  # Directory already exists

# Allowed file extensions
ALLOWED_EXTENSIONS = {
    'ppt', 'pptx', 'doc', 'docx', 'xls', 'xlsx', 'pdf'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    """Serve the main page"""
    return send_from_directory('.', 'index.html')

@app.route('/<path:path>')
def serve_static(path):
    """Serve static files from the current directory"""
    return send_from_directory('.', path)

@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files:
        return 'No file part', 400
    
    file = request.files['file']
    if file.filename == '':
        return 'No selected file', 400
    
    if not allowed_file(file.filename):
        return 'Invalid file type', 400
    
    # Create a temporary directory outside the context manager
    # so it doesn't get deleted immediately
    temp_dir = tempfile.mkdtemp()
    
    try:
        # Save the uploaded file
        input_filename = secure_filename(file.filename)
        input_path = os.path.join(temp_dir, input_filename)
        file.save(input_path)
        logger.info(f"Saved uploaded file to {input_path}")
        
        # Check if this is a PDF to JPG conversion request
        if input_filename.lower().endswith('.pdf') and PDF2IMAGE_AVAILABLE and 'to_jpg' in request.args:
            # Convert PDF to JPG
            output_dir = os.path.join(temp_dir, 'jpg_output')
            os.makedirs(output_dir, exist_ok=True)
            
            # Convert PDF to JPG images
            images = pdf2image.convert_from_path(input_path, dpi=300, fmt="jpeg")
            
            # Save each page as JPG
            jpg_paths = []
            for i, image in enumerate(images):
                jpg_path = os.path.join(output_dir, f"page_{i+1}.jpg")
                image.save(jpg_path, "JPEG")
                jpg_paths.append(jpg_path)
            
            # Create a ZIP file with all the JPGs
            base_name = os.path.splitext(input_filename)[0]
            zip_filename = f"{base_name}_jpg.zip"
            zip_path = os.path.join(temp_dir, zip_filename)
            
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for jpg_path in jpg_paths:
                    zipf.write(jpg_path, os.path.basename(jpg_path))
            
            # Make a copy of the output file to avoid file access issues
            output_copy_dir = tempfile.mkdtemp()
            output_copy_path = os.path.join(output_copy_dir, zip_filename)
            shutil.copy2(zip_path, output_copy_path)
            
            # Send the ZIP file
            response = send_file(
                output_copy_path,
                mimetype='application/zip',
                as_attachment=True,
                download_name=zip_filename
            )
            
            # Set a callback to clean up resources after the response is sent
            @response.call_on_close
            def cleanup():
                try:
                    time.sleep(0.5)  # Small delay to ensure file handles are released
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir, ignore_errors=True)
                    if os.path.exists(output_copy_dir):
                        shutil.rmtree(output_copy_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(f"Cleanup error (non-critical): {str(e)}")
            
            return response
            
        # Check if this is a PDF to PPTX conversion request
        elif input_filename.lower().endswith('.pdf') and 'to_pptx' in request.args:
            # Convert PDF to PPTX
            output_filename = os.path.splitext(input_filename)[0] + '.pptx'
            output_path = os.path.join(temp_dir, output_filename)
            
            # Try to convert using Docker first
            conversion_success = convert_pdf_to_pptx(input_path, output_path, input_filename, output_filename)
            
            if not conversion_success or not os.path.exists(output_path):
                raise Exception("Failed to convert PDF to PPTX. Please try again.")
            
            # Make a copy of the output file to avoid file access issues
            output_copy_dir = tempfile.mkdtemp()
            output_copy_path = os.path.join(output_copy_dir, output_filename)
            shutil.copy2(output_path, output_copy_path)
            
            # Send the converted file
            response = send_file(
                output_copy_path,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=output_filename
            )
            
            # Set a callback to clean up resources after the response is sent
            @response.call_on_close
            def cleanup():
                try:
                    time.sleep(0.5)  # Small delay to ensure file handles are released
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir, ignore_errors=True)
                    if os.path.exists(output_copy_dir):
                        shutil.rmtree(output_copy_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(f"Cleanup error (non-critical): {str(e)}")
            
            return response
        
        # For all file types, convert to PDF using LibreOffice
        output_filename = os.path.splitext(input_filename)[0] + '.pdf'
        output_path = os.path.join(temp_dir, output_filename)
        
        # Create a copy of the final output file in a safe location
        output_copy_dir = tempfile.mkdtemp()
        output_copy_path = os.path.join(output_copy_dir, output_filename)
        
        # Try to convert using Docker first
        conversion_success = convert_using_docker(input_path, output_path, input_filename, output_filename)
        
        # If Docker conversion failed, try local LibreOffice
        if not conversion_success:
            conversion_success = convert_using_local(input_path, temp_dir)
            # Update output_path as local conversion might create a file with different name
            output_path = os.path.join(temp_dir, output_filename)
        
        # Check if conversion was successful
        if not conversion_success or not os.path.exists(output_path):
            raise Exception("Failed to convert file to PDF. Please try again.")
        
        # Make a copy of the output file to avoid file access issues
        try:
            shutil.copy2(output_path, output_copy_path)
            logger.info(f"Created safe copy of output file: {output_copy_path}")
        except Exception as e:
            logger.error(f"Error creating file copy: {str(e)}")
            # If we can't copy, we'll try to use the original with a retry mechanism
            output_copy_path = output_path
        
        # Send the converted file
        try:
            # Return the file in a way that doesn't keep it open on the server
            response = send_file(
                output_copy_path,
                mimetype='application/pdf',
                as_attachment=True,
                download_name=output_filename
            )
            # Set a callback to clean up resources after the response is sent
            @response.call_on_close
            def cleanup():
                try:
                    time.sleep(0.5)  # Small delay to ensure file handles are released
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir, ignore_errors=True)
                    if os.path.exists(output_copy_dir):
                        shutil.rmtree(output_copy_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(f"Cleanup error (non-critical): {str(e)}")
            
            return response
        except Exception as e:
            logger.error(f"Error sending file: {str(e)}")
            raise
    
    except Exception as e:
        logger.error(f"Error during conversion: {str(e)}")
        # Clean up in case of error
        try:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
        except:
            pass
        return str(e), 500

def convert_pdf_to_pptx(input_path, output_path, input_filename, output_filename):
    """Convert PDF to PPTX using Docker LibreOffice with direct method"""
    try:
        logger.info("Attempting PDF to PPTX conversion with Docker")
        
        # Check if Docker container is running
        result = subprocess.run(
            ["docker", "ps", "--filter", f"name={LIBREOFFICE_CONTAINER}"],
            capture_output=True, text=True, check=False
        )
        
        # If container not found in 'docker ps', try to start it
        if LIBREOFFICE_CONTAINER not in result.stdout:
            logger.info("Starting Docker container")
            subprocess.run(["docker", "start", LIBREOFFICE_CONTAINER], 
                           check=False, capture_output=True)
            # Give it a moment to start
            time.sleep(2)
        
        # Create a completely isolated workspace in the container using a random name
        import random
        import string
        random_id = ''.join(random.choices(string.ascii_lowercase + string.digits, k=8))
        container_work_dir = f"/tmp/pdf2pptx_{random_id}"
        
        # Create the directory in the container
        logger.info(f"Creating work directory in container: {container_work_dir}")
        mkdir_cmd = ["docker", "exec", LIBREOFFICE_CONTAINER, "mkdir", "-p", container_work_dir]
        mkdir_result = subprocess.run(mkdir_cmd, capture_output=True, text=True, check=False)
        if mkdir_result.returncode != 0:
            logger.error(f"Failed to create work directory: {mkdir_result.stderr}")
            return False
            
        # Copy PDF file to the container work directory
        container_pdf_path = f"{container_work_dir}/input.pdf"
        logger.info(f"Copying PDF file to container: {container_pdf_path}")
        copy_pdf_cmd = ["docker", "cp", input_path, f"{LIBREOFFICE_CONTAINER}:{container_pdf_path}"]
        copy_result = subprocess.run(copy_pdf_cmd, capture_output=True, text=True, check=False)
        if copy_result.returncode != 0:
            logger.error(f"Failed to copy PDF to container: {copy_result.stderr}")
            return False
            
        # We'll try multiple methods to convert the PDF to PPTX
        
        # Method 1: Direct PDF to PPTX conversion
        logger.info("Trying method 1: Direct PDF to PPTX conversion")
        direct_cmd = [
            "docker", "exec", LIBREOFFICE_CONTAINER,
            "libreoffice", "--headless", "--convert-to", "pptx", 
            "--outdir", container_work_dir, container_pdf_path
        ]
        direct_result = subprocess.run(direct_cmd, capture_output=True, text=True, check=False)
        logger.info(f"Method 1 output: {direct_result.stdout}")
        if direct_result.stderr:
            logger.error(f"Method 1 error: {direct_result.stderr}")
            
        # Method 2: PDF to ODG then to PPTX
        logger.info("Trying method 2: PDF to ODG then to PPTX")
        # First convert to ODG
        container_odg_path = f"{container_work_dir}/temp.odg"
        odg_cmd = [
            "docker", "exec", LIBREOFFICE_CONTAINER,
            "libreoffice", "--headless", "--convert-to", "odg", 
            "--infilter=draw_pdf_import", "--outdir", container_work_dir, 
            container_pdf_path
        ]
        odg_result = subprocess.run(odg_cmd, capture_output=True, text=True, check=False)
        logger.info(f"ODG conversion output: {odg_result.stdout}")
        if odg_result.stderr:
            logger.error(f"ODG conversion error: {odg_result.stderr}")
            
        # Then convert ODG to PPTX
        odg_to_pptx_cmd = [
            "docker", "exec", LIBREOFFICE_CONTAINER,
            "libreoffice", "--headless", "--convert-to", "pptx", 
            "--outdir", container_work_dir, 
            f"{container_work_dir}/input.odg"  # The output of the ODG conversion
        ]
        odg_to_pptx_result = subprocess.run(odg_to_pptx_cmd, capture_output=True, text=True, check=False)
        logger.info(f"ODG to PPTX output: {odg_to_pptx_result.stdout}")
        if odg_to_pptx_result.stderr:
            logger.error(f"ODG to PPTX error: {odg_to_pptx_result.stderr}")
            
        # Method 3: PDF import filter with writer
        logger.info("Trying method 3: Writer PDF import filter")
        writer_cmd = [
            "docker", "exec", LIBREOFFICE_CONTAINER,
            "libreoffice", "--headless", "--infilter=writer_pdf_import", 
            "--convert-to", "pptx", "--outdir", container_work_dir, 
            container_pdf_path
        ]
        writer_result = subprocess.run(writer_cmd, capture_output=True, text=True, check=False)
        logger.info(f"Writer method output: {writer_result.stdout}")
        if writer_result.stderr:
            logger.error(f"Writer method error: {writer_result.stderr}")
            
        # List all files in the work directory to see what was created
        list_cmd = ["docker", "exec", LIBREOFFICE_CONTAINER, "ls", "-la", container_work_dir]
        list_result = subprocess.run(list_cmd, capture_output=True, text=True, check=False)
        logger.info(f"Files in work directory: {list_result.stdout}")
        
        # Look for any PPTX file in the work directory
        find_pptx_cmd = ["docker", "exec", LIBREOFFICE_CONTAINER, "find", container_work_dir, "-name", "*.pptx"]
        find_result = subprocess.run(find_pptx_cmd, capture_output=True, text=True, check=False)
        pptx_files = find_result.stdout.strip().split('\n') if find_result.stdout else []
        logger.info(f"Found PPTX files: {pptx_files}")
        
        if pptx_files and pptx_files[0]:
            # Copy the first PPTX file back
            logger.info(f"Copying PPTX back: {pptx_files[0]} -> {output_path}")
            copy_back_cmd = ["docker", "cp", f"{LIBREOFFICE_CONTAINER}:{pptx_files[0]}", output_path]
            copy_back_result = subprocess.run(copy_back_cmd, capture_output=True, text=True, check=False)
            
            if copy_back_result.returncode != 0:
                logger.error(f"Failed to copy PPTX back: {copy_back_result.stderr}")
                return False
                
            # Verify the file exists and has content
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info("Successfully converted PDF to PPTX")
                
                # Clean up
                cleanup_cmd = ["docker", "exec", LIBREOFFICE_CONTAINER, "rm", "-rf", container_work_dir]
                subprocess.run(cleanup_cmd, capture_output=True, check=False)
                
                return True
        
        # If we got here, all methods failed or the file couldn't be copied back
        logger.error("All conversion methods failed - no valid PPTX generated")
        
        # As a last resort, try to create an empty PPTX with a note about the failure
        try:
            from pptx import Presentation
            
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "PDF Conversion Failed"
            subtitle.text = "The PDF could not be converted to PowerPoint format. Please try a different PDF file."
            
            prs.save(output_path)
            logger.info("Created fallback PPTX with error message")
            return True
        except:
            logger.error("Failed to create fallback PPTX")
            
        return False
        
    except Exception as e:
        logger.error(f"Error in Docker PDF to PPTX conversion: {e}", exc_info=True)
        return False

def convert_using_docker(input_path, output_path, input_filename, output_filename):
    """Convert file using Docker LibreOffice"""
    try:
        logger.info("Attempting Docker conversion")
        
        # Check if Docker container is running
        result = subprocess.run(
            ["docker", "ps", "--filter", f"name={LIBREOFFICE_CONTAINER}"],
            capture_output=True, text=True, check=False
        )
        
        # If container not found in 'docker ps', try to start it
        if LIBREOFFICE_CONTAINER not in result.stdout:
            logger.info("Starting Docker container")
            subprocess.run(["docker", "start", LIBREOFFICE_CONTAINER], 
                           check=False, capture_output=True)
            # Give it a moment to start
            time.sleep(2)
        
        # Copy file to container
        logger.info(f"Copying file to container: {input_filename}")
        copy_to_cmd = ["docker", "cp", input_path, f"{LIBREOFFICE_CONTAINER}:/tmp/{input_filename}"]
        subprocess.run(copy_to_cmd, check=True)
        
        # Convert file in container
        logger.info(f"Converting file in container")
        convert_cmd = [
            "docker", "exec", LIBREOFFICE_CONTAINER,
            "libreoffice", "--headless", "--convert-to", "pdf", 
            "--outdir", "/tmp", f"/tmp/{input_filename}"
        ]
        subprocess.run(convert_cmd, check=True)
        
        # Copy result back
        logger.info(f"Copying result back from container")
        copy_back_cmd = ["docker", "cp", f"{LIBREOFFICE_CONTAINER}:/tmp/{output_filename}", output_path]
        subprocess.run(copy_back_cmd, check=True)
        
        # Verify the file exists
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            logger.info("Docker conversion successful")
            # Clean up files in container
            try:
                cleanup_cmd = [
                    "docker", "exec", LIBREOFFICE_CONTAINER,
                    "rm", "-f", f"/tmp/{input_filename}", f"/tmp/{output_filename}"
                ]
                subprocess.run(cleanup_cmd, check=False)
            except:
                pass
            return True
        else:
            logger.warning("Docker conversion failed - output file not found or empty")
            return False
            
    except Exception as e:
        logger.error(f"Error in Docker conversion: {e}")
        return False

def convert_using_local(input_path, output_dir):
    """Convert file using local LibreOffice installation"""
    try:
        logger.info("Attempting local LibreOffice conversion")
        
        # Find LibreOffice executable
        soffice_path = find_libreoffice()
        if not soffice_path:
            logger.error("LibreOffice not found")
            return False
        
        # Convert file
        convert_cmd = [
            soffice_path, "--headless", "--convert-to", "pdf",
            "--outdir", output_dir, input_path
        ]
        result = subprocess.run(convert_cmd, capture_output=True, text=True, check=False)
        
        if result.returncode != 0:
            logger.error(f"Local conversion failed: {result.stderr}")
            return False
        
        logger.info("Local conversion succeeded")
        return True
    
    except Exception as e:
        logger.error(f"Error in local conversion: {e}")
        return False

def find_libreoffice():
    """Find LibreOffice executable on the system"""
    # Common locations for LibreOffice executable
    possible_paths = [
        # Windows paths
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # Linux paths
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        # macOS paths
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            return path
    
    return None

@app.route('/health', methods=['GET'])
def health_check():
    """Endpoint to check if the server is running"""
    return "Server is running", 200

@app.route('/pdf-to-jpg-convert', methods=['POST'])
def pdf_to_jpg_convert():
    """Special route specifically for PDF to JPG conversion"""
    if not PDF2IMAGE_AVAILABLE:
        return "PDF to JPG conversion is not available. Please install pdf2image library.", 503
    
    # Redirect to the main convert route with a special parameter
    return convert_file()

@app.route('/pdf-to-pptx-convert', methods=['POST'])
def pdf_to_pptx_convert():
    """Special route specifically for PDF to PPTX conversion"""
    # Redirect to the main convert route with a special parameter
    request.args = {'to_pptx': 'true'}
    return convert_file()

if __name__ == "__main__":
    if not PDF2IMAGE_AVAILABLE:
        logger.warning("pdf2image is not installed. PDF to JPG conversion will not be available.")
    app.run(host="0.0.0.0", port=5000, debug=True) 